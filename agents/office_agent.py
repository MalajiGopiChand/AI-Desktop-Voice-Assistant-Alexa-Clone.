"""Office automation — Word, Excel, PowerPoint, PDF."""
import os
from agents.base_agent import BaseAgent
from core.llm_client import chat, FAST_MODEL


class OfficeAgent(BaseAgent):
    def __init__(self):
        super().__init__("office_agent")

    def execute(self, action, params):
        try:
            handlers = {
                "read_word": self._read_word,
                "write_word": self._write_word,
                "read_excel": self._read_excel,
                "summarize_excel": self._summarize_excel,
                "read_pdf": self._read_pdf,
                "summarize_pdf": self._summarize_pdf,
                "create_report": self._create_report,
                "read_pptx": self._read_pptx,
            }
            handler = handlers.get(action)
            if not handler:
                return {"success": False, "message": f"Unknown action: {action}", "data": {}}
            return handler(params)
        except Exception as e:
            return {"success": False, "message": str(e), "data": {}}

    def _read_word(self, params):
        try:
            from docx import Document
        except ImportError:
            return {"success": False, "message": "Install python-docx: pip install python-docx", "data": {}}
        path = os.path.expanduser(params.get("path", ""))
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return {"success": True, "message": text[:500], "data": {"text": text}}

    def _write_word(self, params):
        from docx import Document
        path = params.get("path", "report.docx")
        content = params.get("content", "")
        doc = Document()
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(path)
        return {"success": True, "message": f"Document saved to {path}", "data": {"path": path}}

    def _read_excel(self, params):
        import openpyxl
        path = os.path.expanduser(params.get("path", ""))
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheet = wb.active
        rows = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= 20:
                break
            rows.append([str(c) if c is not None else "" for c in row])
        wb.close()
        preview = "\n".join(", ".join(r) for r in rows[:10])
        return {"success": True, "message": preview, "data": {"rows": rows}}

    def _summarize_excel(self, params):
        result = self._read_excel(params)
        text = result["data"].get("rows", [])
        flat = str(text)[:4000]
        summary = chat(
            [{"role": "system", "content": "Summarize this spreadsheet data briefly."},
             {"role": "user", "content": flat}],
            model=FAST_MODEL, max_tokens=250,
        )
        return {"success": True, "message": summary, "data": {"summary": summary}}

    def _read_pdf(self, params):
        import pdfplumber
        path = os.path.expanduser(params.get("path", ""))
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages[:10]:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        text = "\n".join(text_parts)
        return {"success": True, "message": text[:500], "data": {"text": text}}

    def _summarize_pdf(self, params):
        result = self._read_pdf(params)
        text = result["data"].get("text", "")
        summary = chat(
            [{"role": "system", "content": "Summarize this document for voice output."},
             {"role": "user", "content": text[:6000]}],
            model=FAST_MODEL, max_tokens=300,
        )
        return {"success": True, "message": summary, "data": {"summary": summary}}

    def _read_pptx(self, params):
        from pptx import Presentation
        path = os.path.expanduser(params.get("path", ""))
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides.append(shape.text.strip())
        text = "\n".join(slides)
        return {"success": True, "message": text[:500], "data": {"text": text}}

    def _create_report(self, params):
        topic = params.get("topic", "Report")
        content = chat(
            [{"role": "system", "content": "Write a professional report with sections."},
             {"role": "user", "content": f"Create a report about: {topic}"}],
            model=FAST_MODEL, max_tokens=800,
        )
        from config import REPORTS_DIR
        path = os.path.join(REPORTS_DIR, f"report_{topic.replace(' ', '_')[:30]}.docx")
        return self._write_word({"path": path, "content": content})
